import io
import os
from PIL import Image, ImageEnhance
import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv

dir_path = 'static/presentations'

load_dotenv()

# Set your Adobe Stock API key
adobe_stock_api_key = "Set your Adobe Stock API key"


def search_adobe_stock_images(query):
    headers = {
        'x-api-key': adobe_stock_api_key,
        'x-product': 'myApp'
    }
    response = requests.get(
        f'https://stock.adobe.io/Rest/Media/1/Search/Files?locale=en_US&search_parameters[words]={query}',
        headers=headers)

    if response.status_code == 200:
        data = response.json()
        if 'files' in data and len(data['files']) > 0:
            return data['files'][0]['thumbnail_url']
        else:
            print("No images found for the query.")
            return None
    else:
        print(f"Error: {response.status_code} - {response.text}")
        return None


def parse_response(response):
    slides = response.split('\n\n')
    slides_content = []
    for slide in slides:
        lines = slide.split('\n')
        title_line = lines[0]
        if ': ' in title_line:
            title = title_line.split(': ', 1)[1]  # Extract the title after 'Slide X: '
        else:
            title = title_line
        content_lines = []
        keywords = ""
        image_search_query = ""  # Initialize empty image search query
        for line in lines[1:]:
            if line.startswith('Content:'):
                continue
            elif line.startswith('Keywords:'):
                keywords = line.split(': ', 1)[1] if ': ' in line else ""
            elif line.startswith('Image Search Query:'):
                image_search_query = line.split(': ', 1)[1] if ': ' in line else ""  # Capture image search query
            else:
                content_lines.append(line)
        content = '\n'.join(content_lines)  # Join the lines to form the content
        slides_content.append({'title': title, 'content': content, 'keywords': keywords})
    return slides_content


def delete_first_two_slides(presentation):
    slide_ids = [1, 0]
    for slide_id in slide_ids:
        if slide_id < len(presentation.slides):
            xml_slides = presentation.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[slide_id])


def add_logo_with_opacity(slide, logo_path, position, width, opacity=0.5):
    # Load the image
    logo_image = Image.open(logo_path)

    # Apply opacity (if needed)
    if opacity < 1.0:
        alpha = logo_image.split()[3]
        alpha = ImageEnhance.Brightness(alpha).enhance(opacity)
        logo_image.putalpha(alpha)

    # Save the modified logo image to a BytesIO object
    image_stream = io.BytesIO()
    logo_image.save(image_stream, format='PNG')
    image_stream.seek(0)

    # Add the image to the slide
    slide.shapes.add_picture(image_stream, *position, width=width)


def create_ppt(slides_content, template_choice, presentation_title, presenter_name, insert_image):
    # Load the selected template
    template_path = os.path.join(dir_path, f"{template_choice}.pptx")
    logo_path = os.path.join('static', 'images', 'logo.png')

    prs = Presentation(template_path)

    # Define layouts and other settings
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]
    table_of_contents_layout = prs.slide_layouts[1]

    logo_position = (Inches(0.5), prs.slide_height - Inches(1.2))
    logo_width = Inches(1.0)

    # Add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = presentation_title
    subtitle = slide.placeholders[1]
    subtitle.text = f"Presented by {presenter_name}"

    # Customize title slide based on template
    if template_choice == 'Mustard_template':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0, 0, 0)
        for paragraph in subtitle.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0, 0, 0)

    add_logo_with_opacity(slide, logo_path, position=logo_position, width=logo_width, opacity=0.5)

    # Add table of contents slide
    index_slide = prs.slides.add_slide(table_of_contents_layout)
    title_placeholder = index_slide.shapes.title
    title_placeholder.text = "Table of Contents"
    for paragraph in title_placeholder.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(20)
            run.font.bold = True
            if template_choice == "dark_modern":
                run.font.color.rgb = RGBColor(255, 165, 0)
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)

    content_placeholder = index_slide.placeholders[1]
    content_text_frame = content_placeholder.text_frame
    content_text_frame.clear()
    for i, slide_content in enumerate(slides_content):
        p = content_text_frame.add_paragraph()
        p.text = f"{i + 1}. {slide_content['title']}"
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.LEFT
        if template_choice == 'dark_modern':
            p.font.name = 'Times New Roman'
            p.font.color.rgb = RGBColor(255, 165, 0)
        elif template_choice == 'bright_modern':
            p.font.name = 'Arial'
            p.font.color.rgb = RGBColor(255, 20, 147)
        elif template_choice == 'geometric_template':
            p.font.name = 'Comic Sans MS'
            p.font.color.rgb = RGBColor(0, 0, 0)
        elif template_choice == 'Mustard_template':
            p.font.name = 'Calibri'
            p.font.color.rgb = RGBColor(0, 0, 0)

    add_logo_with_opacity(index_slide, logo_path, position=logo_position, width=logo_width, opacity=0.5)

    # Add content slides with template-specific image placement
    for slide_content in slides_content:
        slide = prs.slides.add_slide(content_slide_layout)

        # Only add the title and content, skipping the keyword and image generation prompt
        title_text = slide_content['title']
        content_text = slide_content['content']

        if template_choice == 'geometric_template':
            # Specific adjustments for geometric_template
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 1:
                    placeholder.text = title_text
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Comic Sans MS'
                            run.font.color.rgb = RGBColor(0, 0, 0)
                elif placeholder.placeholder_format.type == 7:
                    placeholder.text = content_text
                    placeholder.text_frame.paragraphs[0].font.size = Pt(20)

            # Image placement for geometric_template
            if insert_image:
                try:
                    image_url = search_adobe_stock_images(slide_content['keywords'])
                    if image_url:
                        image_data = requests.get(image_url).content
                        image_stream = io.BytesIO(image_data)
                        image_width = Inches(3.5)
                        image_height = Inches(2.5)
                        left = Inches(9.2)
                        top = Inches(4.2)
                        slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)
                except Exception as e:
                    print(f"Failed to add image: {e}")

        elif template_choice == 'Mustard_template':
            # Specific adjustments for Mustard_template
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 1:
                    placeholder.text = title_text
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Calibri'
                            run.font.color.rgb = RGBColor(0, 0, 0)
                elif placeholder.placeholder_format.type == 7:
                    placeholder.text = content_text
                    placeholder.text_frame.paragraphs[0].font.size = Pt(19)
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(0, 0, 0)

            # Image placement for Mustard_template
            if insert_image:
                try:
                    image_url = search_adobe_stock_images(slide_content['keywords'])
                    if image_url:
                        image_data = requests.get(image_url).content
                        image_stream = io.BytesIO(image_data)
                        image_width = Inches(3.5)
                        image_height = Inches(2.5)
                        left = Inches(6.2)
                        top = Inches(4.7)
                        slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)
                except Exception as e:
                    print(f"Failed to add image: {e}")

        elif template_choice == 'fun_template':
            # Specific adjustments for fun_template
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 1:
                    placeholder.text = title_text
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Comic Sans MS'
                            run.font.color.rgb = RGBColor(0, 0, 0)
                elif placeholder.placeholder_format.type == 7:
                    placeholder.text = content_text
                    placeholder.text_frame.paragraphs[0].font.size = Pt(20)

            # Image placement for fun_template
            if insert_image:
                try:
                    image_url = search_adobe_stock_images(slide_content['keywords'])
                    if image_url:
                        image_data = requests.get(image_url).content
                        image_stream = io.BytesIO(image_data)
                        image_width = Inches(3.5)
                        image_height = Inches(2.5)
                        left = Inches(8.5)
                        top = Inches(3.5)
                        slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)
                except Exception as e:
                    print(f"Failed to add image: {e}")

        else:
            # General settings for other templates
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 1:
                    placeholder.text = title_text
                elif placeholder.placeholder_format.type == 7:
                    placeholder.text = content_text

            # General image placement
            if insert_image:
                try:
                    image_url = search_adobe_stock_images(slide_content['keywords'])
                    if image_url:
                        image_data = requests.get(image_url).content
                        image_stream = io.BytesIO(image_data)
                        image_width = Inches(8)
                        image_height = Inches(5)
                        left = prs.slide_width - image_width
                        top = prs.slide_height - image_height - Inches(4)
                        slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)
                except Exception as e:
                    print(f"Failed to add image: {e}")

        # Add the logo to each content slide with 50% opacity
        add_logo_with_opacity(slide, logo_path, position=logo_position, width=logo_width, opacity=0.5)

    # Add thank you slide
    thank_you_slide_layout = prs.slide_layouts[1]
    thank_you_slide = prs.slides.add_slide(thank_you_slide_layout)
    thank_you_title = thank_you_slide.shapes.title
    thank_you_title.text = "Thank You"

    if template_choice == 'Mustard_template':
        for paragraph in thank_you_title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0, 0, 0)

    # Add the logo to the thank you slide with 50% opacity
    add_logo_with_opacity(thank_you_slide, logo_path, position=logo_position, width=logo_width, opacity=0.5)

    # Delete the first two slides after all new slides have been added
    delete_first_two_slides(prs)

    # Save the presentation after everything is processed
    prs.save(os.path.join('generated', 'generated_presentation.pptx'))
    print("Presentation created and saved successfully.")
